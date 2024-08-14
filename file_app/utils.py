from io import BytesIO

from docx import Document
from django.conf import settings
from pptx import Presentation
from predibase import PredibaseClient
from PyPDF2 import PdfReader
from rest_framework import exceptions


def get_summary_from_predibase(content):
    """
    Get a summary for the content using Predibase API.
    """
    client = PredibaseClient(token=settings.PREDIBASE_API_KEY)
    llm_deployment = client.LLM("pb://deployments/llama-3-1-8b-instruct")
    data = {
        "text": f"Summarize the following text in 3 lines:\n\n{content}\n\nSummary:"
    }
    llm_response = llm_deployment.prompt(
        data=data,
        temperature=0.7,
        max_new_tokens=150,
        bypass_system_prompt=False
    )
    return llm_response.response


def get_file_summary(file):
    """
    Get a summary for the uploaded file by parsing the file first and then hitting the predibase API
    """
    try:
        if file.name.split('.')[-1] == "pdf":
            content = read_pdf_file(file)
        elif file.name.split('.')[-1] == "docx":
            content = read_docx_file(file)
        else:
            content = read_pptx_file(file)
        return get_summary_from_predibase(content)
    except Exception as e:
        raise exceptions.ValidationError({'error': 'Error occurred while reading file', 'extra': str(e)}) from e


def read_pdf_file(file):
    """
    Get the text content of the pdf file
    """
    pdf_reader = PdfReader(BytesIO(file.read()))
    file.seek(0)
    text_content = ""
    for page in pdf_reader.pages:
        text_content += page.extract_text()
    return text_content

def read_docx_file(file):
    """
    Get the text content of the docx file
    """
    doc = Document(BytesIO(file.read()))
    file.seek(0)
    full_text = []
    # Extract text from paragraphs
    for para in doc.paragraphs:
        full_text.append(para.text)
    # Extract text from tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    full_text.append(paragraph.text)
    # Extract text from text boxes and other shapes
    for shape in doc.inline_shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                full_text.append(paragraph.text)
    return '\n'.join(filter(None, full_text))

def read_pptx_file(file):
    """
    Get the text content of the pptx file
    """
    prs = Presentation(BytesIO(file.read()))
    file.seek(0)
    text_content = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text_content.append(shape.text)
    return '\n'.join(filter(None, text_content))
